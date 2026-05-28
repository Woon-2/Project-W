from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "poster_output" / "Hi-Z_Visibility_Animation_Poster_A0_source_captures.pptx"
OUT = ROOT / "poster_output" / "Hi-Z_Visibility_Animation_Poster_A0_source_captures_with_conclusion.pptx"

FONT_KR = "Malgun Gothic"
NAVY = RGBColor(21, 45, 73)
PURPLE = RGBColor(94, 74, 145)
GRAY = RGBColor(51, 58, 67)
MUTED = RGBColor(100, 110, 121)
LINE = RGBColor(188, 199, 210)
WHITE = RGBColor(255, 255, 255)
REF_BG = RGBColor(238, 242, 246)
REF_LINE = RGBColor(224, 230, 236)


def set_run(run, size, color=GRAY, bold=False):
    run.font.name = FONT_KR
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def set_text(shape, text, size, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size, color, bold)


def add_text(slide, x, y, w, h, text, size=14, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    set_text(shape, text, size, color, bold, align)
    return shape


def add_bullets(slide, x, y, w, h, items, size=12.7):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.0)
    tf.margin_bottom = Cm(0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(2.4)
        p.line_spacing = 1.02
        r = p.add_run()
        r.text = f"- {item}"
        set_run(r, size, GRAY)
    return shape


def move_and_shrink_references(slide):
    ref_texts = {
        "참고문헌": (4.0, 107.25, 75.8, 0.45, 8.0, True),
        "[1] B. Dudash": (4.0, 107.85, 75.8, 1.2, 6.6, False),
        "[4] U. Haar": (4.0, 109.0, 75.8, 0.75, 6.6, False),
    }
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        for prefix, (x, y, w, h, size, bold) in ref_texts.items():
            if text.startswith(prefix):
                shape.left = Cm(x)
                shape.top = Cm(y)
                shape.width = Cm(w)
                shape.height = Cm(h)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        set_run(run, size, MUTED, bold)

    # Move the quiet reference background rectangle if present.
    for shape in slide.shapes:
        if (
            abs(shape.left.cm - 3.2) < 0.2
            and abs(shape.top.cm - 98.0) < 0.4
            and abs(shape.width.cm - 77.7) < 0.5
            and abs(shape.height.cm - 8.3) < 0.7
        ):
            shape.top = Cm(106.85)
            shape.height = Cm(3.25)
            shape.fill.solid()
            shape.fill.fore_color.rgb = REF_BG
            shape.line.color.rgb = REF_LINE


def add_conclusion(slide):
    x, y, w, h = 3.2, 97.25, 77.7, 8.9
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LINE
    panel.line.width = Pt(1)
    panel.adjustments[0] = 0.035

    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(1.65))
    head.fill.solid()
    head.fill.fore_color.rgb = PURPLE
    head.line.color.rgb = PURPLE
    head.text_frame.clear()
    head.text_frame.margin_left = Cm(0.55)
    head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = head.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "결론"
    set_run(r, 21, WHITE, True)

    add_text(slide, x + 0.8, y + 2.35, 10.0, 0.8, "연구의 기여", 14.5, NAVY, True)
    add_bullets(
        slide,
        x + 0.8,
        y + 3.25,
        36.2,
        4.3,
        [
            "GPU-CPU 동기화 없이 가능한 최신 가시성 정보를 CPU 애니메이션 시스템에 반영하는 피드백 구조를 제시하였다.",
            "GPU-Driven 렌더링 파이프라인과 연계 가능하며, 하드웨어 인스턴싱 단위를 유지하면서 적용할 수 있다.",
            "GPU 가시성 정보를 애니메이션 스케줄링에 반영해 복잡한 CPU 애니메이션 연산을 줄일 수 있음을 보였다.",
        ],
        12.3,
    )

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + 39.0), Cm(y + 2.35), Cm(0.03), Cm(5.45))
    divider.fill.solid()
    divider.fill.fore_color.rgb = LINE
    divider.line.color.rgb = LINE

    add_text(slide, x + 40.1, y + 2.35, 10.0, 0.8, "향후 연구", 14.5, NAVY, True)
    add_bullets(
        slide,
        x + 40.1,
        y + 3.25,
        36.5,
        4.3,
        [
            "삼각형 클러스터 단위 GPU-Driven 파이프라인에서 더 세밀한 가시성 역전파를 연구할 수 있다.",
            "가시성 피드백을 물리 엔진, 전투 시스템 등 CPU 중심 게임 로직 최적화로 확장할 수 있다.",
        ],
        12.3,
    )


def main():
    prs = Presentation(SRC)
    slide = prs.slides[0]
    move_and_shrink_references(slide)
    add_conclusion(slide)
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
