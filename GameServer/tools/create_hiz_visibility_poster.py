from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
PDF_FILE = ROOT.parent / "Hi-Z 오클루전 컬링 결과를 활용한 가시성 기반 애니메이션 갱신 기법.pdf"
OUT_DIR = ROOT / "poster_output"
CAPTURE_DIR = OUT_DIR / "paper_captures"
OUT_FILE = OUT_DIR / "Hi-Z_Visibility_Animation_Poster_A0_source_captures.pptx"

FONT_KR = "Malgun Gothic"

NAVY = RGBColor(21, 45, 73)
BLUE = RGBColor(30, 91, 150)
TEAL = RGBColor(0, 119, 124)
GREEN = RGBColor(43, 130, 79)
ORANGE = RGBColor(204, 112, 39)
GRAY = RGBColor(51, 58, 67)
MUTED = RGBColor(100, 110, 121)
LINE = RGBColor(188, 199, 210)
BG = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(229, 239, 249)
PALE_TEAL = RGBColor(226, 244, 241)
PALE_ORANGE = RGBColor(252, 238, 224)


def crop_pdf_assets():
    """Crop source paper figures/tables for poster use."""
    CAPTURE_DIR.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(PDF_FILE)
    page = doc[1]
    # Coordinates are PDF points on page 2. Crops include original captions.
    crops = {
        "fig1_system.png": fitz.Rect(35, 55, 370, 196),
        "table1_visibility_cases.png": fitz.Rect(224, 205, 389, 297),
        "fig2_missing_instance.png": fitz.Rect(247, 388, 389, 449),
        "table2_performance.png": fitz.Rect(475, 41, 640, 96),
    }
    paths = {}
    matrix = fitz.Matrix(3.0, 3.0)
    for name, rect in crops.items():
        pix = page.get_pixmap(matrix=matrix, clip=rect, alpha=False)
        path = CAPTURE_DIR / name
        pix.save(path)
        paths[name] = path
    return paths


def set_run(run, size, color=GRAY, bold=False):
    run.font.name = FONT_KR
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, x, y, w, h, text, size=18, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.12)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = text
    set_run(r, size, color, bold)
    return shape


def add_panel(slide, x, y, w, h, title, color=BLUE):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = LINE
    panel.line.width = Pt(1)
    panel.adjustments[0] = 0.035

    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(2.2))
    head.fill.solid()
    head.fill.fore_color.rgb = color
    head.line.color.rgb = color
    head.text_frame.clear()
    head.text_frame.margin_left = Cm(0.55)
    head.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = head.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, 26, WHITE, True)
    return panel


def add_bullets(slide, x, y, w, h, items, size=17.5, color=GRAY, gap=5):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.space_after = Pt(gap)
        p.line_spacing = 1.04
        r = p.add_run()
        r.text = f"- {item}"
        set_run(r, size, color)
    return shape


def add_tag(slide, x, y, w, h, label, value, fill, accent):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = Pt(1.1)
    box.adjustments[0] = 0.08
    add_text(slide, x + 0.25, y + 0.25, w - 0.5, 0.65, label, 12.5, accent, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.25, y + 0.95, w - 0.5, h - 1.05, value, 17.5, NAVY, True, PP_ALIGN.CENTER)


def add_picture_with_caption(slide, path, x, y, w, caption=None):
    pic = slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    pic.line.color.rgb = LINE
    pic.line.width = Pt(0.75)
    if caption:
        add_text(slide, x, y + pic.height.cm + 0.1, w, 0.55, caption, 10.5, MUTED, False, PP_ALIGN.CENTER)
    return pic


def build():
    captures = crop_pdf_assets()

    prs = Presentation()
    prs.slide_width = Cm(84.1)
    prs.slide_height = Cm(118.9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Header
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(84.1), Cm(12.6))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY
    add_text(slide, 3.0, 1.45, 78.1, 4.0, "Hi-Z 오클루전 컬링 결과를 활용한\n가시성 기반 애니메이션 갱신 기법", 34, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 3.0, 6.25, 78.1, 1.1, "A Visibility-based Animation Update Method Using Hi-Z Occlusion Culling", 18, RGBColor(221, 232, 242), False, PP_ALIGN.CENTER)
    add_text(slide, 3.0, 8.25, 78.1, 1.0, "장명운○, 우상훈, 이종진, 정내훈 | 한국공학대학교 게임공학과", 17.5, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, 3.0, 9.8, 78.1, 0.9, "wkdauddns9@gmail.com, asm1133@naver.com, ljj3640@gmail.com, nhjung@tukorea.ac.kr", 12.5, RGBColor(221, 232, 242), False, PP_ALIGN.CENTER)

    # Overview
    add_panel(slide, 3.2, 14.0, 77.7, 11.2, "개요", BLUE)
    add_text(
        slide,
        4.2,
        16.9,
        49.0,
        4.2,
        "GPU에서 판정된 Hi-Z 기반 가시성 정보를 CPU 애니메이션 시스템에 비동기적으로 반영하여, 화면에 기여하지 않는 객체의 애니메이션 갱신을 선택적으로 생략한다.",
        22,
        NAVY,
        True,
    )
    add_tag(slide, 55.0, 16.1, 7.2, 4.4, "갱신 시간", "4.3ms -> 2.1ms", PALE_TEAL, GREEN)
    add_tag(slide, 63.3, 16.1, 7.2, 4.4, "피드백", "2,416 byte", PALE_BLUE, BLUE)
    add_tag(slide, 71.6, 16.1, 7.2, 4.4, "병합 비용", "< 0.1ms", PALE_ORANGE, ORANGE)
    add_text(slide, 4.2, 21.5, 74.7, 1.5, "핵심 효과: 애니메이션 연산량을 전체 객체 수가 아니라 실제 가시 객체 수에 비례하도록 낮춘다.", 17, GRAY)

    # Purpose
    add_panel(slide, 3.2, 27.0, 26.0, 28.7, "연구 목적", TEAL)
    add_bullets(
        slide,
        4.2,
        30.0,
        24.0,
        10.4,
        [
            "렌더링 최적화에 머물던 GPU 가시성 정보를 CPU 애니메이션 갱신 최적화에 활용한다.",
            "CPU-GPU 동기화로 인한 파이프라인 지연 없이 가시성 피드백을 얻는다.",
            "IK, 래그돌 물리, 실시간 보간처럼 CPU 비용이 큰 애니메이션 환경에서 병목을 완화한다.",
        ],
        17.2,
    )
    add_picture_with_caption(slide, captures["table2_performance.png"], 5.0, 42.0, 20.2, "원문 [표 2] 성능 분석표")
    add_text(slide, 4.2, 49.8, 23.8, 3.1, "성능 평가는 Ryzen 7 7800X3D, RTX 4070 Ti SUPER, Windows 11 환경에서 수행되었다.", 15.5, MUTED)

    # Core method, large middle/right region
    add_panel(slide, 31.0, 27.0, 49.9, 68.7, "가시성 기반 애니메이션 갱신", NAVY)
    add_picture_with_caption(slide, captures["fig1_system.png"], 32.2, 30.2, 28.0, "원문 [그림 1] 시스템 개요")
    add_text(
        slide,
        61.2,
        30.4,
        18.0,
        11.2,
        "렌더링 이벤트는 메시-재질 단위로 정렬되고, Hi-Z 오클루전 컬링은 각 이벤트에 대한 32비트 가시성 플래그를 생성한다. 하위 비트는 가시성, 상위 비트는 객체 식별자를 저장한다.",
        15.7,
        GRAY,
    )
    add_picture_with_caption(slide, captures["table1_visibility_cases.png"], 32.4, 48.0, 17.2, "원문 [표 1] 가시성 판정 결과 조합")
    add_picture_with_caption(slide, captures["fig2_missing_instance.png"], 51.5, 48.0, 15.7, "원문 [그림 2] 인스턴스 누락 상황")
    add_text(
        slide,
        68.3,
        47.9,
        10.6,
        15.1,
        "동기화 없이 버퍼를 읽으면 이전/현재 프레임의 가시성 정보가 섞일 수 있다. 특히 인스턴스 재배치 시 실제로 보이던 객체가 비가시로 잘못 병합되는 누락 상황이 발생할 수 있다.",
        14.8,
        GRAY,
    )

    # Dense explanation blocks inside core section
    add_tag(slide, 32.4, 64.3, 14.7, 4.7, "1. 가시성 플래그", "visible bit + object ID", PALE_BLUE, BLUE)
    add_tag(slide, 48.2, 64.3, 14.7, 4.7, "2. 링 버퍼", "2-slot visibility buffer", PALE_ORANGE, ORANGE)
    add_tag(slide, 64.0, 64.3, 14.7, 4.7, "3. 객체 병합", "OR 기반 가시 판정", PALE_TEAL, GREEN)
    add_bullets(
        slide,
        32.4,
        71.0,
        22.5,
        12.0,
        [
            "CPU는 GPU와 명시적으로 동기화하지 않고 가시성 플래그 버퍼를 읽는다.",
            "객체 단위 가시성은 매 프레임 비가시로 초기화한 뒤, 구성 인스턴스 중 하나라도 보이면 가시로 갱신한다.",
            "슬롯별 컬링 대상 수를 따로 저장해 유효 범위 밖의 오래된 데이터를 읽지 않는다.",
        ],
        15.7,
        GRAY,
        3,
    )
    add_bullets(
        slide,
        56.0,
        71.0,
        22.5,
        12.0,
        [
            "애니메이션 시스템은 가시 객체만 대상으로 거리와 마지막 갱신 시점을 고려해 우선순위를 계산한다.",
            "우선순위 큐에서 제한 시간 안에 가능한 만큼만 갱신하고 나머지는 생략한다.",
            "결과적으로 컬링과 LOD가 렌더링뿐 아니라 애니메이션 갱신에도 적용된다.",
        ],
        15.7,
        GRAY,
        3,
    )
    formula = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(33.0), Cm(84.8), Cm(45.8), Cm(8.0))
    formula.fill.solid()
    formula.fill.fore_color.rgb = RGBColor(249, 251, 253)
    formula.line.color.rgb = LINE
    formula.adjustments[0] = 0.04
    add_text(slide, 34.1, 85.7, 15.0, 1.0, "애니메이션 우선순위", 16, NAVY, True)
    add_text(slide, 34.1, 87.4, 43.4, 2.9, "w_d = 1 / (1 + (d x k_d)^2)     w_t = 1 + t / (t + k_t)     priority = w_d x w_t", 17.5, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, 34.1, 90.6, 43.4, 1.1, "가까운 객체와 오래 갱신되지 않은 객체의 갱신 우선순위를 높인다.", 14.5, MUTED, False, PP_ALIGN.CENTER)

    # Result
    add_panel(slide, 3.2, 57.5, 26.0, 38.2, "결과", GREEN)
    add_text(slide, 4.2, 60.5, 23.8, 3.2, "Hi-Z 오클루전 컬링과 시야 절두체 컬링을 모두 반영했을 때 평균 애니메이션 갱신 시간은 2.1ms였다.", 17.2, GRAY)
    add_text(slide, 4.2, 64.3, 23.8, 2.8, "시야 절두체 컬링만 반영한 경우는 4.3ms로, 제안 기법은 약 2배 수준의 성능 향상을 보였다.", 17.2, GRAY)
    add_tag(slide, 4.4, 68.0, 11.0, 5.5, "Baseline", "4.3 ms", PALE_ORANGE, ORANGE)
    add_tag(slide, 16.4, 68.0, 11.0, 5.5, "Proposed", "2.1 ms", PALE_TEAL, GREEN)
    add_tag(slide, 4.4, 75.0, 11.0, 5.5, "Readback", "2,416 B", PALE_BLUE, BLUE)
    add_tag(slide, 16.4, 75.0, 11.0, 5.5, "CPU cost", "< 0.1 ms", PALE_TEAL, GREEN)
    add_text(
        slide,
        4.2,
        82.4,
        23.8,
        7.9,
        "해석: readback과 객체 가시성 병합 비용이 매우 작아, CPU 기반 애니메이션이 주요 병목이 되는 장면에서 효과가 크다. 렌더링 품질을 유지하면서 보이지 않는 객체의 애니메이션 계산을 줄이는 것이 핵심이다.",
        16.2,
        GRAY,
    )

    # References, intentionally quiet
    ref = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3.2), Cm(98.0), Cm(77.7), Cm(8.3))
    ref.fill.solid()
    ref.fill.fore_color.rgb = RGBColor(238, 242, 246)
    ref.line.color.rgb = RGBColor(224, 230, 236)
    add_text(slide, 4.0, 98.5, 75.8, 0.8, "참고문헌", 11.5, MUTED, True)
    add_text(
        slide,
        4.0,
        99.6,
        75.8,
        3.9,
        "[1] B. Dudash, Skinned Instancing, NVIDIA SDK 10, 2007.  [2] N. Greene et al., Hierarchical Z-buffer visibility, SIGGRAPH, 1993.  [3] H. Zhang et al., Visibility Culling Using Hierarchical Occlusion Maps, SIGGRAPH, 1997.",
        9.4,
        MUTED,
    )
    add_text(
        slide,
        4.0,
        103.2,
        75.8,
        1.7,
        "[4] U. Haar and S. Aaltonen, GPU-Driven Rendering Pipelines, SIGGRAPH Course, 2015.  [5] B. Karis et al., Nanite, SIGGRAPH Course, 2021.",
        9.4,
        MUTED,
    )

    # Footer
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(111.6), Cm(84.1), Cm(7.3))
    foot.fill.solid()
    foot.fill.fore_color.rgb = NAVY
    foot.line.color.rgb = NAVY
    add_text(slide, 3.0, 112.5, 58.0, 1.6, "본 연구는 2026년도 과학기술정보통신부 및 정보통신기획평가원의 'SW중심대학사업' 지원을 받아 수행되었음(2025-0-00050)", 12.5, WHITE)
    add_text(slide, 61.5, 112.5, 19.5, 1.3, "2026 한국게임학회 춘계 학술발표대회", 12.5, RGBColor(221, 232, 242), False, PP_ALIGN.RIGHT)

    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE

    OUT_DIR.mkdir(exist_ok=True)
    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    print(build())
